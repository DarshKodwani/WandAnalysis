"""Generate a PowerPoint version of the three-ecosystem workflow diagram."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Slide dimensions: widescreen landscape ---
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(22)

slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)

# --- Colours ---
COL_STEP = RGBColor(0x2C, 0x3E, 0x50)
COL_FSL = RGBColor(0x29, 0x80, 0xB9)
COL_SPM = RGBColor(0xE7, 0x4C, 0x3C)
COL_FMRIPREP = RGBColor(0x27, 0xAE, 0x60)
COL_FSL_BG = RGBColor(0xD6, 0xEA, 0xF8)
COL_SPM_BG = RGBColor(0xFA, 0xDB, 0xD8)
COL_FMRIPREP_BG = RGBColor(0xD5, 0xF5, 0xE3)
COL_GREY_BG = RGBColor(0xF8, 0xF9, 0xF9)
COL_ARROW = RGBColor(0x7F, 0x8C, 0x8D)
COL_BLACK = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# --- Data ---
steps = [
    "1. Motion Correction\n(MC)",
    "2. Slice-Timing\nCorrection (STC)",
    "3. Distortion\nCorrection (SDC)",
    "4. Spatial\nNormalisation (NORM)",
    "5. Spatial\nSmoothing (SMOOTH)",
    "6. Temporal\nFiltering (FILT)",
    "7. Nuisance\nRegression (REG)",
    "8. Parcellation\n(PARC)",
    "9. Connectivity\nMetric (METRIC)",
]

fsl_tools = [
    "MCFLIRT\nref=mid, normcorr",
    "slicetimer\nsinc, ref=mid slice",
    "FUGUE (3T)\nTOPUP (7T)",
    "FLIRT + FNIRT\n→ MNI152 2mm",
    "SUSAN\n6 mm FWHM",
    "Butterworth\n0.01–0.1 Hz",
    "8P\n(6P + WM/CSF)",
    "Schaefer 200\nmean extraction",
    "Pearson\n+ Fisher-z",
]

spm_tools = [
    "SPM Realign\nref=first, LSQ",
    "SPM STC\nFourier, ref=mid",
    "SPM FieldMap\n/ VDM",
    "DARTEL\n→ MNI152",
    "Gaussian\n8 mm FWHM",
    "Butterworth\n0.01–0.1 Hz",
    "24P Friston\n+ WM/CSF",
    "AAL3\nmean extraction",
    "Pearson\n+ Fisher-z",
]

fmriprep_tools = [
    "MCFLIRT\nref=mid, normcorr",
    "Fourier\nref=0.5 (mid-TR)",
    "TOPUP (7T)\nSyN-SDC (3T)",
    "ANTs SyN\n→ MNI152NLin2009c",
    "None\n(0 mm)",
    "None\n(user applies)",
    "aCompCor\n+ ICA-AROMA",
    "Schaefer 200\nmean extraction",
    "Pearson\n+ Fisher-z",
]


def add_rounded_box(slide, left, top, width, height, fill_color, border_color, text, font_size, font_color, bold=False):
    """Add a rounded rectangle with centred text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    # Adjust corner rounding
    shape.adjustments[0] = 0.1

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Vertical centring
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

    # Split text by newlines into separate paragraphs
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bold

    # Vertical centre
    shape.text_frame.paragraphs[0].space_before = Pt(2)
    shape.text_frame.word_wrap = True
    # Use anchor middle
    from pptx.oxml.ns import qn
    txBody = shape.text_frame._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    bodyPr.set("anchor", "ctr")

    return shape


def add_arrow(slide, x, y_start, y_end):
    """Add a downward arrow connector."""
    # Use a simple line shape
    connector = slide.shapes.add_connector(
        1,  # straight connector
        x, y_start, x, y_end)
    connector.line.color.rgb = COL_ARROW
    connector.line.width = Pt(1.5)
    # Add arrowhead at end
    from pptx.oxml.ns import qn
    ln = connector.line._ln
    tail = ln.makeelement(qn("a:tailEnd"), {})
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    ln.append(tail)


# --- Layout constants (in inches) ---
box_w = Inches(3.0)
box_h = Inches(1.4)
spacing_y = Inches(1.9)  # vertical distance between row centres

x_step = Inches(0.5)
x_fsl = Inches(4.0)
x_spm = Inches(7.5)
x_fmriprep = Inches(11.0)

y_title = Inches(0.2)
y_headers = Inches(1.0)
y_input = Inches(1.6)
y_start = Inches(2.8)

# --- Title ---
txBox = slide.shapes.add_textbox(Inches(0), y_title, Inches(16), Inches(0.7))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Three-Ecosystem Default Pipelines"
run.font.size = Pt(24)
run.font.bold = True
run.font.color.rgb = COL_STEP

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "9-step path from raw BOLD to FC matrix"
run2.font.size = Pt(14)
run2.font.italic = True
run2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

# --- Column headers ---
for xc, label, color in [
    (x_fsl, "FSL", COL_FSL),
    (x_spm, "SPM", COL_SPM),
    (x_fmriprep, "fMRIPrep", COL_FMRIPREP),
]:
    txBox = slide.shapes.add_textbox(xc, y_headers, box_w, Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = color

# --- Input box ---
add_rounded_box(slide, x_step, y_input, box_w, Inches(0.6),
                COL_GREY_BG, COL_STEP, "Raw 4-D BOLD", 12, COL_STEP, bold=True)

# --- Step rows ---
for i in range(9):
    y = y_start + i * spacing_y

    # Step label
    add_rounded_box(slide, x_step, y, box_w, box_h,
                    COL_GREY_BG, COL_STEP, steps[i], 11, COL_STEP, bold=True)

    # FSL
    add_rounded_box(slide, x_fsl, y, box_w, box_h,
                    COL_FSL_BG, COL_FSL, fsl_tools[i], 10, COL_BLACK)

    # SPM
    add_rounded_box(slide, x_spm, y, box_w, box_h,
                    COL_SPM_BG, COL_SPM, spm_tools[i], 10, COL_BLACK)

    # fMRIPrep
    add_rounded_box(slide, x_fmriprep, y, box_w, box_h,
                    COL_FMRIPREP_BG, COL_FMRIPREP, fmriprep_tools[i], 10, COL_BLACK)

    # Arrows to next row
    if i < 8:
        y_arrow_start = y + box_h
        y_arrow_end = y + spacing_y
        arrow_mid = x_step + box_w / 2
        for xc in [x_step, x_fsl, x_spm, x_fmriprep]:
            add_arrow(slide, xc + box_w / 2, y_arrow_start, y_arrow_end)

# --- Output row ---
y_output = y_start + 8 * spacing_y + box_h + Inches(0.3)

# Arrows from last step row to output
for xc in [x_step, x_fsl, x_spm, x_fmriprep]:
    add_arrow(slide, xc + box_w / 2, y_start + 8 * spacing_y + box_h, y_output)

add_rounded_box(slide, x_step, y_output, box_w, Inches(0.6),
                COL_GREY_BG, COL_STEP, "P × P FC Matrix", 12, COL_STEP, bold=True)

for xc, label, bg, ec in [
    (x_fsl, "FSL.DEFAULT\nFC Matrix", COL_FSL_BG, COL_FSL),
    (x_spm, "SPM.DEFAULT\nFC Matrix", COL_SPM_BG, COL_SPM),
    (x_fmriprep, "FMRIPREP.DEFAULT\nFC Matrix", COL_FMRIPREP_BG, COL_FMRIPREP),
]:
    add_rounded_box(slide, xc, y_output, box_w, Inches(0.6),
                    bg, ec, label, 10, COL_BLACK, bold=True)


# --- Save ---
outpath = "/MRIWork/MRIWork06/nr/darsh_kodwani/WandAnalysis/docs/paper/when_is_a_result_robust_marginalization_over_reasonable_analysis_choices_in_fmri/workflow_diagram.pptx"
prs.save(outpath)
print(f"Saved: {outpath}")
